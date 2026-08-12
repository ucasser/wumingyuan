"""文档解析：把 md / pdf / epub 读成纯文本 + 结构信息。

统一返回 List[Section]，每个 Section 是文档的一个自然片段：
  - md   : 按标题层级切成 section，heading 记录标题路径
  - pdf  : 每页一个 section，heading 记为 "p.N"
  - epub : 每个章节一个 section，heading 记为章节标题
分块（chunker）在此基础上再做定长切分。
"""
import os
from dataclasses import dataclass, field
from typing import List


@dataclass
class Section:
    text: str
    heading: str = ""            # 标题 / 页码 / 章节名
    order: int = 0               # 在文档内的顺序


@dataclass
class Document:
    path: str
    title: str
    sections: List[Section] = field(default_factory=list)


# 纯文本类（含 markdown）：走 _parse_md，有 # 标题时按层级切，否则整篇一段
TEXT_EXTS = {".md", ".markdown", ".txt", ".text", ".rst", ".org", ".tex", ".log"}


def parse(path: str, ocr: dict = None) -> Document:
    ext = os.path.splitext(path)[1].lower()
    if ext in TEXT_EXTS:
        return _parse_md(path)
    if ext == ".pdf":
        return _parse_pdf(path, ocr=ocr)
    if ext == ".epub":
        return _parse_epub(path)
    if ext == ".docx":
        return _parse_docx(path)
    if ext == ".pptx":
        return _parse_pptx(path)
    if ext == ".rtf":
        return _parse_rtf(path)
    if ext in (".html", ".htm"):
        return _parse_html(path)
    if ext == ".doc":
        raise ValueError("旧版 .doc（二进制）不支持，请在 Word 里另存为 .docx 后再导入")
    if ext == ".ppt":
        raise ValueError("旧版 .ppt（二进制）不支持，请在 PowerPoint 里另存为 .pptx 后再导入")
    raise ValueError(f"不支持的文件类型: {ext}")


# ---------------- Markdown ----------------
def _parse_md(path: str) -> Document:
    with open(path, "r", encoding="utf-8", errors="ignore") as f:
        lines = f.readlines()

    title = os.path.splitext(os.path.basename(path))[0]
    sections: List[Section] = []
    stack = []            # 当前标题栈 [(level, text)]
    buf = []
    order = 0

    def flush():
        nonlocal buf, order
        text = "".join(buf).strip()
        if text:
            heading = " > ".join(t for _, t in stack)
            sections.append(Section(text=text, heading=heading, order=order))
            order += 1
        buf = []

    for ln in lines:
        stripped = ln.lstrip()
        if stripped.startswith("#"):
            level = len(stripped) - len(stripped.lstrip("#"))
            htext = stripped[level:].strip()
            if 1 <= level <= 6 and htext:
                flush()
                while stack and stack[-1][0] >= level:
                    stack.pop()
                stack.append((level, htext))
                if level == 1 and not sections:
                    title = htext
                continue
        buf.append(ln)
    flush()

    if not sections:  # 空文件兜底
        sections.append(Section(text="", heading="", order=0))
    return Document(path=path, title=title, sections=sections)


# ---------------- PDF ----------------
_OCR_WARNED = False


def _ocr_page(page, ocr: dict) -> str:
    """对单页做 OCR。缺 Tesseract 时优雅返回空串（只提示一次）。"""
    global _OCR_WARNED
    try:
        import io
        import pytesseract
        from PIL import Image
        from . import ocr as ocr_mod
        ocr_mod.configure_pytesseract()      # Windows 下把 tesseract 路径喂给 pytesseract
        pix = page.get_pixmap(dpi=int(ocr.get("dpi", 200)))
        img = Image.open(io.BytesIO(pix.tobytes("png")))
        return pytesseract.image_to_string(img, lang=ocr.get("lang", "chi_sim+eng")).strip()
    except Exception as ex:
        if not _OCR_WARNED:
            print(f"[ocr] OCR 不可用（需系统安装 Tesseract 及语言包）：{ex}")
            _OCR_WARNED = True
        return ""


def _parse_pdf(path: str, ocr: dict = None) -> Document:
    import fitz  # PyMuPDF

    doc = fitz.open(path)
    file_stem = os.path.splitext(os.path.basename(path))[0]
    meta_title = (doc.metadata.get("title") or "").strip()
    
    # 清理无意义的 PDF 元数据伪标题（如 !00001, NONE, 纯数字, 临时word名）
    if not meta_title or re.search(r"^[!0-9\s]+$|^NONE$|^Microsoft Word|^Untitled", meta_title, re.I):
        title = file_stem
    else:
        title = meta_title
    ocr_on = bool(ocr and ocr.get("enabled"))
    min_chars = int((ocr or {}).get("min_chars", 20))
    page_texts = [page.get_text("text").strip() for page in doc]
    low_pages = [i + 1 for i, text in enumerate(page_texts) if len(text) < min_chars]

    # 飞桨模式：只把缺少文本层的页一次性提交云端，返回逐页 Markdown。
    # Markdown 仍是普通文本，后续继续走原来的本地 chunker + bge-m3。
    paddle_texts = {}
    provider = (ocr or {}).get("provider", "tesseract")
    if ocr_on and provider == "paddle" and low_pages:
        try:
            from . import paddle_ocr
            base_dir = os.path.dirname(os.path.dirname(__file__))
            paddle_texts = paddle_ocr.parse_pdf(
                path, (ocr or {}).get("paddle", {}), low_pages, base_dir
            )
        except Exception as ex:
            if not (ocr or {}).get("fallback_local", True):
                doc.close()
                raise RuntimeError(f"飞桨 OCR 失败：{ex}") from ex
            print(f"[ocr] 飞桨 OCR 失败，回退本地 Tesseract：{ex}")

    sections = []
    for i, page in enumerate(doc):
        text = page_texts[i]
        if ocr_on and len(text) < min_chars:
            if provider == "paddle" and paddle_texts.get(i + 1):
                ocr_text = paddle_texts[i + 1]
            elif provider == "tesseract" or (ocr or {}).get("fallback_local", True):
                ocr_text = _ocr_page(page, ocr or {})
            else:
                ocr_text = ""
            if len(ocr_text) > len(text):
                text = ocr_text
        if text:
            sections.append(Section(text=text, heading=f"p.{i + 1}", order=i))
    doc.close()
    if not sections:
        sections.append(Section(text="", heading="p.1", order=0))
    return Document(path=path, title=title, sections=sections)


# ---------------- EPUB ----------------
def _parse_epub(path: str) -> Document:
    from ebooklib import epub, ITEM_DOCUMENT
    from bs4 import BeautifulSoup

    book = epub.read_epub(path)
    md = book.get_metadata("DC", "title")
    title = md[0][0] if md else os.path.splitext(os.path.basename(path))[0]

    sections = []
    order = 0
    for item in book.get_items_of_type(ITEM_DOCUMENT):
        soup = BeautifulSoup(item.get_content(), "lxml")
        for tag in soup(["script", "style"]):
            tag.decompose()
        text = soup.get_text("\n").strip()
        if not text:
            continue
        h = soup.find(["h1", "h2", "h3"])
        heading = h.get_text(strip=True) if h else f"section {order + 1}"
        sections.append(Section(text=text, heading=heading, order=order))
        order += 1
    if not sections:
        sections.append(Section(text="", heading="", order=0))
    return Document(path=path, title=title, sections=sections)


# ---------------- Word (.docx) ----------------
def _parse_docx(path: str) -> Document:
    from docx import Document as Docx

    d = Docx(path)
    title = (d.core_properties.title or "").strip() \
        or os.path.splitext(os.path.basename(path))[0]

    sections: List[Section] = []
    stack = []          # 标题栈 [(level, text)]
    buf = []
    order = 0

    def flush():
        nonlocal buf, order
        text = "\n".join(buf).strip()
        if text:
            heading = " > ".join(t for _, t in stack)
            sections.append(Section(text=text, heading=heading, order=order))
            order += 1
        buf = []

    for para in d.paragraphs:
        text = para.text.strip()
        style = (para.style.name if para.style else "") or ""
        if style.startswith("Heading") and text:
            # 样式名形如 "Heading 1"，取末尾数字作层级
            digits = "".join(c for c in style if c.isdigit())
            level = int(digits) if digits else 1
            flush()
            while stack and stack[-1][0] >= level:
                stack.pop()
            stack.append((level, text))
            if level == 1 and not sections and title == \
                    os.path.splitext(os.path.basename(path))[0]:
                title = text
            continue
        if text:
            buf.append(text)
    flush()

    # 表格：逐行拼成 "单元格 | 单元格"
    for ti, table in enumerate(d.tables):
        rows = []
        for row in table.rows:
            cells = [c.text.strip() for c in row.cells]
            if any(cells):
                rows.append(" | ".join(cells))
        txt = "\n".join(rows).strip()
        if txt:
            sections.append(Section(text=txt, heading=f"表格{ti + 1}", order=order))
            order += 1

    if not sections:
        sections.append(Section(text="", heading="", order=0))
    return Document(path=path, title=title, sections=sections)


# ---------------- PowerPoint (.pptx) ----------------
def _parse_pptx(path: str) -> Document:
    from pptx import Presentation

    prs = Presentation(path)
    title = os.path.splitext(os.path.basename(path))[0]
    sections = []
    for i, slide in enumerate(prs.slides):
        texts = []
        for shape in slide.shapes:
            if getattr(shape, "has_text_frame", False):
                t = shape.text_frame.text.strip()
                if t:
                    texts.append(t)
            if getattr(shape, "has_table", False):
                for row in shape.table.rows:
                    cells = [c.text.strip() for c in row.cells]
                    if any(cells):
                        texts.append(" | ".join(cells))
        # 演讲者备注
        if slide.has_notes_slide:
            note = slide.notes_slide.notes_text_frame.text.strip()
            if note:
                texts.append("备注：" + note)
        body = "\n".join(texts).strip()
        if body:
            sections.append(Section(text=body, heading=f"幻灯片{i + 1}", order=i))
    if not sections:
        sections.append(Section(text="", heading="", order=0))
    return Document(path=path, title=title, sections=sections)


# ---------------- RTF ----------------
def _parse_rtf(path: str) -> Document:
    from striprtf.striprtf import rtf_to_text

    with open(path, "r", encoding="utf-8", errors="ignore") as f:
        raw = f.read()
    text = rtf_to_text(raw).strip()
    title = os.path.splitext(os.path.basename(path))[0]
    return Document(path=path, title=title,
                    sections=[Section(text=text, heading="", order=0)])


# ---------------- HTML ----------------
def _parse_html(path: str) -> Document:
    from bs4 import BeautifulSoup

    with open(path, "r", encoding="utf-8", errors="ignore") as f:
        soup = BeautifulSoup(f.read(), "lxml")
    for tag in soup(["script", "style", "noscript"]):
        tag.decompose()
    title = (soup.title.get_text(strip=True) if soup.title else "") \
        or os.path.splitext(os.path.basename(path))[0]
    text = soup.get_text("\n").strip()
    # 压掉连续空行
    text = "\n".join(ln for ln in (l.rstrip() for l in text.splitlines()) if ln)
    return Document(path=path, title=title,
                    sections=[Section(text=text, heading="", order=0)])
